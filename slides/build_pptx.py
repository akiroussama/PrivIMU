"""Build slides/PrivIMU_final.pptx from the rendered PDF deck.

Strategy: each PDF page becomes a high-DPI PNG and is embedded as a full-bleed
slide background in a 16:9 PowerPoint (13.333" x 7.5"). The slide title and
speaker notes are added as PPT metadata so the deck is navigable in PowerPoint
and presenter view shows the speaker script.
"""

from __future__ import annotations

import json
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DECK_PDF = ROOT / "slides" / "PrivIMU_final.pdf"
OUT_PPTX = ROOT / "slides" / "PrivIMU_final.pptx"
RENDER_DIR = ROOT / "slides" / "_pptx_pages"
SPEAKER_DIR = ROOT / "slides" / "speaker_notes"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SLIDE_TITLES: list[str] = [
    "PrivIMU — Cover",
    "IMU sensors are everywhere in IoT",
    "Anonymous ≠ non-identifiable",
    "Threat model: from anonymous stream to top-3 identity",
    "From activity recognition to identity leakage",
    "MotionSense — 24 subjects, 6 activities, iPhone in pocket",
    "A reproducible privacy pipeline",
    "How the signal becomes evidence",
    "Evaluation: avoid inflated results",
    "Live demo: from CSV to privacy risk",
    "A single window already narrows the search",
    "Blue team: how loud must we be to hide?",
    "Results — measured, not hand-written",
    "Limits & privacy-by-design mitigations",
    "Anonymous ≠ non-identifiable. Treat motion data as PII.",
]

# Map each slide (1-indexed) to the speaker-note source file.
SPEAKER_MAP = {
    1: "M1_intro.md", 2: "M1_intro.md", 3: "M1_intro.md",
    4: "M2_sota.md", 5: "M2_sota.md",
    6: "M3_methodo.md", 7: "M3_methodo.md", 8: "M3_methodo.md", 9: "M3_methodo.md",
    10: "M4_demo.md", 11: "M4_demo.md", 12: "M4_demo.md",
    13: "M5_results.md", 14: "M5_results.md", 15: "M5_results.md",
}


def render_pages(dpi: int = 200) -> list[Path]:
    RENDER_DIR.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(DECK_PDF))
    out: list[Path] = []
    for i, page in enumerate(doc, start=1):
        pix = page.get_pixmap(dpi=dpi, alpha=False)
        target = RENDER_DIR / f"slide_{i:02d}.png"
        pix.save(str(target))
        out.append(target)
    doc.close()
    return out


def load_speaker_note(slide_idx: int) -> str:
    fname = SPEAKER_MAP.get(slide_idx)
    if not fname:
        return ""
    p = SPEAKER_DIR / fname
    if not p.exists():
        return ""
    return p.read_text(encoding="utf-8")


def build(images: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # blank
    for idx, img_path in enumerate(images, start=1):
        slide = prs.slides.add_slide(blank_layout)
        # Full-bleed background image
        slide.shapes.add_picture(
            str(img_path), 0, 0, width=SLIDE_W, height=SLIDE_H
        )
        # Speaker notes
        notes_tf = slide.notes_slide.notes_text_frame
        title = SLIDE_TITLES[idx - 1] if idx - 1 < len(SLIDE_TITLES) else f"Slide {idx}"
        speaker = load_speaker_note(idx)
        notes_tf.text = f"[Slide {idx} — {title}]\n\n{speaker}".strip()

    prs.save(str(OUT_PPTX))


def main() -> None:
    if not DECK_PDF.exists():
        raise SystemExit(f"PDF deck not found: {DECK_PDF}. Run Chrome print first.")
    images = render_pages(dpi=200)
    build(images)
    summary = {
        "pptx": str(OUT_PPTX.relative_to(ROOT)),
        "n_slides": len(images),
        "size_mb": round(OUT_PPTX.stat().st_size / 1_048_576, 3),
    }
    print(json.dumps(summary, indent=2))


if __name__ == "__main__":
    main()
